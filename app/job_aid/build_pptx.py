"""Render a structured job-aid spec into a copy of the HR Job Aid Template.

Design notes (learned the hard way in the agent):
- Screencast slides 6/7/8 are LAYOUT OPTIONS. We pick one and duplicate it for
  every screen group, then delete the unused layouts.
- Duplicated slides must keep their white <p:bg> override, or they fall back to
  the dark "Segue Midnight" layout background and the dark text disappears.
- Picture slots are always re-filled with freshly added images (real video
  frames when available, else a labeled placeholder) so image relationships are
  valid on the duplicated slide parts.
"""
from __future__ import annotations

import copy
import io
import math
import os
import tempfile
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

DARK_TEXT = RGBColor(0x01, 0x18, 0x2E)

# Template slide indices (0-based) for the three screencast layouts.
LAYOUT_SLIDES = {4: "decision", 5: "screencast4", 6: "screencast2", 7: "screencast2b"}
SCREENCAST_IDXS = [5, 6, 7]


def render(spec: dict, template_path: str, frame_paths: Optional[list[str]] = None) -> bytes:
    """Return .pptx bytes for the filled job aid."""
    prs = Presentation(template_path)
    frame_paths = frame_paths or []

    _fill_title(prs.slides[0], spec)
    _fill_disclaimer(prs.slides[1], spec)
    _fill_overview([prs.slides[2], prs.slides[3]], spec)
    _fill_decision(prs.slides[4], spec)
    _build_screencast(prs, spec, frame_paths)
    _fill_resources(prs.slides[-1], spec)  # slide 9 stays last

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- section fillers --------------------------------------------------------

def _fill_title(slide, spec):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if "<INSERT JOB AID TITLE>" in t or "JOB AID TITLE" in t.upper():
            _set_text(shape, spec.get("title", "Job Aid"))
        elif "subtitle" in t.lower():
            _set_text(shape, spec.get("subtitle", ""))


def _fill_disclaimer(slide, spec):
    disclaimers = spec.get("disclaimers", []) or []
    di = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if "<Insert Disclaimer>" in t:
            _set_text(shape, disclaimers[di] if di < len(disclaimers) else "")
            di += 1
        elif "<Insert Document Title>" in t:
            _set_text(shape, spec.get("document_title", spec.get("title", "")))


def _fill_overview(slides, spec):
    cards = spec.get("overview", []) or []
    ci = 0
    for slide in slides:
        topics = _shapes_with(slide, "Topic")  # 'Text Placeholder 1' holds "Topic"
        infos = _shapes_with(slide, "Important info")
        # Fall back to name-based lookup if placeholder text was already changed.
        pairs = list(zip(sorted(topics, key=_yx), sorted(infos, key=_yx)))
        for topic_shape, info_shape in pairs:
            if ci >= len(cards):
                break
            _set_text(topic_shape, cards[ci].get("topic", ""))
            _set_text(info_shape, cards[ci].get("info", ""))
            ci += 1
        for shape in slide.shapes:
            if shape.has_text_frame and "<Insert Document Title>" in shape.text_frame.text:
                _set_text(shape, spec.get("document_title", spec.get("title", "")))


def _fill_decision(slide, spec):
    dt = spec.get("decision_tree") or {}
    points = dt.get("points", []) or []
    pi = 0
    for shape in sorted([s for s in slide.shapes if s.has_text_frame], key=_yx):
        t = shape.text_frame.text
        if "<Insert Decision Point Content>" in t:
            _set_text(shape, points[pi] if pi < len(points) else "")
            pi += 1
        elif "<Insert Stop Point Verbiage>" in t:
            _set_text(shape, dt.get("stop", ""))
        elif "<Insert Next Step Verbiage>" in t:
            _set_text(shape, dt.get("next", ""))
        elif "<Insert Notes/Important Message>" in t:
            _set_text(shape, dt.get("notes", spec.get("notes_message", "")))


# Slide-6 quadrant slots, in reading order (TL, TR, BL, BR), paired by the
# template's shape names. Slide 6 is the one layout with top-level pictures; its
# first step's caption/number live inside a group, so we address shapes by name
# (recursively) rather than by geometry.
SLIDE6_SLOTS = [
    ("TextBox 14", "Oval 12", "Picture 46"),  # top-left
    ("TextBox 15", "Oval 13", "Picture 2"),   # top-right
    ("TextBox 18", "Oval 16", "Picture 8"),   # bottom-left
    ("TextBox 19", "Oval 17", "Picture 45"),  # bottom-right
]
PER_SLIDE = 4
SCREENCAST_TEMPLATE_IDX = 5  # template slide 6


def _build_screencast(prs, spec, frame_paths):
    steps = spec.get("steps", []) or []
    if not steps:
        _delete_slides(prs, SCREENCAST_IDXS)
        return
    n_slides = math.ceil(len(steps) / PER_SLIDE)
    id_lst = prs.slides._sldIdLst
    original_els = [list(id_lst)[i] for i in SCREENCAST_IDXS]  # capture by identity

    new_slides = [_duplicate_slide(prs, SCREENCAST_TEMPLATE_IDX) for _ in range(n_slides)]

    step_no = 0
    for si, slide in enumerate(new_slides):
        group = steps[si * PER_SLIDE:(si + 1) * PER_SLIDE]
        # slide title, notes, running header
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tx = shape.text_frame.text
            if tx.strip().startswith("Topic Title") or "Step By Step" in tx:
                _set_text(shape, spec.get("title", ""), color=DARK_TEXT)
            elif "<Insert Notes/Important Message>" in tx:
                _set_text(shape, spec.get("notes_message", ""), color=DARK_TEXT)
            elif "<Insert Document Title>" in tx:
                _set_text(shape, spec.get("document_title", spec.get("title", "")))
        # quadrant slots
        for j, (cap_name, num_name, pic_name) in enumerate(SLIDE6_SLOTS):
            cap = _find(slide, cap_name)
            num = _find(slide, num_name)
            pic = _find(slide, pic_name)
            if j < len(group):
                step_no += 1
                if cap is not None:
                    _set_text(cap, group[j].get("caption", ""), color=DARK_TEXT)
                if num is not None:
                    _set_text(num, str(step_no))
                if pic is not None:
                    rect = (_emu(pic.left), _emu(pic.top), _emu(pic.width), _emu(pic.height))
                    pic._element.getparent().remove(pic._element)
                    img = _image_for(group[j], step_no, frame_paths)
                    slide.shapes.add_picture(img, rect[0], rect[1], width=rect[2], height=rect[3])
            else:  # unused quadrant on the last slide — clear it
                if cap is not None:
                    _set_text(cap, "")
                if num is not None:
                    _set_text(num, "")
                if pic is not None:
                    pic._element.getparent().remove(pic._element)

    _place_after(prs, new_slides, after_idx=4)   # after the decision-tree slide
    for el in original_els:                        # drop the three original layouts
        id_lst.remove(el)


def _fill_resources(slide, spec):
    res = spec.get("resources", []) or []
    lines = [f"{r.get('label', '')}: {r.get('url', '')}".strip(" :") for r in res]
    blocks = sorted([s for s in slide.shapes if s.has_text_frame and "<Insert Content>" in s.text_frame.text], key=_yx)
    for i, shape in enumerate(blocks):
        _set_text(shape, lines[i] if i < len(lines) else "")
    for shape in slide.shapes:
        if shape.has_text_frame and "<Insert Document Title>" in shape.text_frame.text:
            _set_text(shape, spec.get("document_title", spec.get("title", "")))


# --- slide plumbing ---------------------------------------------------------

def _duplicate_slide(prs, index):
    """Deep-copy a slide's shapes + white background onto a new slide."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):  # strip layout-provided placeholders
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    # Preserve the slide-level white background override (critical!).
    src_bg = source.element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if src_bg is not None:
        c_sld = dest.element.find(qn("p:cSld"))
        old = c_sld.find(qn("p:bg"))
        if old is not None:
            c_sld.remove(old)
        c_sld.insert(0, copy.deepcopy(src_bg))
    return dest


def _place_after(prs, slides, after_idx):
    """Reorder the sldIdLst so `slides` (appended at the end) sit after after_idx."""
    id_lst = prs.slides._sldIdLst
    ids = list(id_lst)
    anchor = ids[after_idx]
    moving = [el for el in ids if prs.slides[list(id_lst).index(el)] in slides] if False else None
    # Simpler: match by element identity of the new slides' sldId.
    new_ids = []
    for s in slides:
        for el in list(id_lst):
            if el.get(qn("r:id")) == _rid_of(prs, s):
                new_ids.append(el)
                break
    for el in new_ids:
        id_lst.remove(el)
    pos = list(id_lst).index(anchor) + 1
    for offset, el in enumerate(new_ids):
        id_lst.insert(pos + offset, el)


def _rid_of(prs, slide):
    part = prs.part
    for rid, rel in part.rels.items():
        if rel.reltype.endswith("/slide") and rel.target_part is slide.part:
            return rid
    return None


def _delete_slides(prs, indexes):
    id_lst = prs.slides._sldIdLst
    to_remove = [list(id_lst)[i] for i in sorted(indexes)]
    for el in to_remove:
        id_lst.remove(el)


# --- shape helpers ----------------------------------------------------------

def _find(container, name):
    """Find a shape by name, descending into groups. Returns the shape or None."""
    for s in container.shapes:
        if s.name == name:
            return s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = _find(s, name)
            if found is not None:
                return found
    return None


def _caption_shapes(slide):
    out = []
    for s in slide.shapes:
        if s.shape_type != MSO_SHAPE_TYPE.TEXT_BOX or not s.has_text_frame:
            continue
        t = s.text_frame.text
        if "Action Verb" in t or "Insert Text Description" in t:
            out.append(s)
    return out


def _number_ovals(slide):
    out = []
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.has_text_frame:
            t = s.text_frame.text.strip()
            if t.isdigit() and len(t) <= 2:
                out.append(s)
    return out


def _shapes_with(slide, needle):
    return [s for s in slide.shapes if s.has_text_frame and needle in s.text_frame.text]


def _image_for(step, step_no, frame_paths):
    idx = step_no - 1
    if idx < len(frame_paths) and frame_paths[idx] and os.path.exists(frame_paths[idx]):
        return frame_paths[idx]
    return _placeholder_png(step_no, step.get("caption", ""), step.get("alt_text", ""))


def _placeholder_png(step_no, caption, alt):
    from PIL import Image, ImageDraw

    W, H = 1080, 540
    img = Image.new("RGB", (W, H), (245, 246, 248))
    d = ImageDraw.Draw(img)
    d.rectangle([6, 6, W - 6, H - 6], outline=(20, 24, 46), width=4)
    d.text((40, 40), "SCREENSHOT PLACEHOLDER", fill=(20, 24, 46))
    d.text((40, 90), f"Step {step_no}", fill=(20, 24, 46))
    d.text((40, 140), f"Capture: {(alt or caption)[:80]}", fill=(20, 24, 46))
    fd, path = tempfile.mkstemp(suffix=".png", prefix="jobaid_ph_")
    os.close(fd)
    img.save(path)
    return path


def _set_text(shape, text, color=None):
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for extra in para.runs[1:]:
            extra._r.getparent().remove(extra._r)
        run = para.runs[0]
    else:
        run = para.add_run()
        run.text = text
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    if color is not None:
        run.font.color.rgb = color


def _emu(v):
    return int(v) if v is not None else 0


def _yx(shape):
    return (_emu(shape.top), _emu(shape.left))
