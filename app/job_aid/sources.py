"""Source ingestion: turn an uploaded file or a pasted link into text the model
can reason over, and (for videos) extract real frames to use as screenshots.

Supported inputs:
- Word (.docx)      -> extracted paragraph text
- PowerPoint (.pptx)-> extracted slide text
- Plain text (.txt) -> as-is (e.g. a pasted transcript)
- Video (.mp4/.mov/.mkv/.webm) -> transcript text is NOT auto-derived (ask the
  user for a transcript), but real frames can be extracted at given timestamps
- URL               -> best-effort fetch of readable text; YouTube pages do not
  expose transcripts to a plain fetch, so the user should paste a transcript.

This mirrors the agent's "Source material" rules: the source describes the
process; screenshots are produced by us (frame extraction for videos).
"""
from __future__ import annotations

import os
import tempfile
from dataclasses import dataclass, field
from typing import Optional
from urllib.request import urlopen, Request

VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".webm", ".m4v", ".avi"}


@dataclass
class Source:
    """Normalized view of whatever the user provided."""
    kind: str                      # 'docx' | 'pptx' | 'text' | 'video' | 'url'
    text: str = ""                 # extracted / provided text (may be empty for video)
    video_path: Optional[str] = None  # local path to a video file, if any
    notes: list[str] = field(default_factory=list)  # warnings to surface in the UI


def from_upload(filename: str, data: bytes) -> Source:
    """Build a Source from an uploaded file's name + bytes."""
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".docx":
        return Source(kind="docx", text=_read_docx(data))
    if ext == ".pptx":
        return Source(kind="pptx", text=_read_pptx(data))
    if ext in (".txt", ".md", ".vtt", ".srt"):
        return Source(kind="text", text=data.decode("utf-8", errors="replace"))
    if ext in VIDEO_EXTS:
        path = _spill_to_tmp(data, ext)
        return Source(
            kind="video",
            video_path=path,
            notes=[
                "Video uploaded. Frames will be extracted for screenshots, but the "
                "spoken/on-screen steps are not auto-transcribed — paste a transcript "
                "or step list in the text box for best results."
            ],
        )
    raise ValueError(f"Unsupported file type: {ext}")


def from_url(url: str) -> Source:
    """Best-effort fetch of a URL's readable text."""
    notes = []
    text = ""
    try:
        req = Request(url, headers={"User-Agent": "Mozilla/5.0 (JobAidCreator)"})
        with urlopen(req, timeout=20) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
        text = _strip_html(raw)
    except Exception as exc:  # noqa: BLE001 - surface, don't crash the app
        notes.append(f"Could not fetch the link ({exc}). Paste the content/transcript instead.")
    if "youtube.com" in url or "youtu.be" in url:
        notes.append(
            "YouTube links do not expose a transcript to a server fetch. Paste the "
            "transcript, or upload the video file so frames can be extracted."
        )
    return Source(kind="url", text=text, notes=notes)


def extract_frames(video_path: str, timestamps_s: list[float]) -> list[str]:
    """Extract one frame per timestamp (seconds); return list of PNG file paths.

    Uses OpenCV so no ffmpeg install is required. Frames are written to a temp dir.
    """
    import cv2  # imported lazily so the app starts even if cv2 is missing

    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise RuntimeError("Could not open the video for frame extraction.")
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
    out_dir = tempfile.mkdtemp(prefix="jobaid_frames_")
    paths: list[str] = []
    for i, ts in enumerate(timestamps_s):
        frame_no = int(round(ts * fps))
        if total:
            frame_no = max(0, min(total - 1, frame_no))
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_no)
        ok, frame = cap.read()
        if not ok:
            continue
        out = os.path.join(out_dir, f"frame_{i + 1}.png")
        cv2.imwrite(out, frame)
        paths.append(out)
    cap.release()
    return paths


# --- internal helpers -------------------------------------------------------

def _spill_to_tmp(data: bytes, ext: str) -> str:
    fd, path = tempfile.mkstemp(suffix=ext, prefix="jobaid_src_")
    with os.fdopen(fd, "wb") as fh:
        fh.write(data)
    return path


def _read_docx(data: bytes) -> str:
    import docx  # python-docx

    path = _spill_to_tmp(data, ".docx")
    doc = docx.Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_pptx(data: bytes) -> str:
    from pptx import Presentation

    path = _spill_to_tmp(data, ".pptx")
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return "\n".join(parts)


def _strip_html(html: str) -> str:
    import re

    html = re.sub(r"(?is)<(script|style).*?</\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", html)
    text = re.sub(r"\s+", " ", text)
    return text.strip()[:20000]
